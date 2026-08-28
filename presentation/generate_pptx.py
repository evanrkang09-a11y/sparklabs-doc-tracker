"""
SparkLabs 서류 추적기 — Korean Presentation Generator (v2)
16 slides, 15-20 minutes
Output: sparklabs-presentation-v2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x2B, 0x55)
MID_BLUE    = RGBColor(0x1A, 0x56, 0xAA)
BRIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xFF)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
DARK        = RGBColor(0x1E, 0x29, 0x3B)
GRAY        = RGBColor(0x64, 0x74, 0x8B)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)

def rgb(r, g, b): return RGBColor(r, g, b)

W, H = Inches(13.33), Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line_color=None, line_width=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def txb(slide, text, l, t, w, h, size, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, font="Malgun Gothic", wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box

def bullet_frame(slide, items, l, t, w, h, size=17, color=DARK,
                 bullet_color=BRIGHT_BLUE, line_space=1.2, font="Malgun Gothic"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = para._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_space * 100000)))
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", "100")
        br = para.add_run()
        br.text = "● "
        br.font.size = Pt(size * 0.62)
        br.font.color.rgb = bullet_color
        br.font.name = font
        tr = para.add_run()
        if isinstance(item, tuple):
            tr.text = item[0]
            tr.font.size = Pt(item[1] if len(item) > 1 else size)
            tr.font.bold = item[2] if len(item) > 2 else False
            tr.font.color.rgb = item[3] if len(item) > 3 else color
        else:
            tr.text = item
            tr.font.size = Pt(size)
            tr.font.color.rgb = color
        tr.font.name = font
    return box

def header_bar(slide, title, sub=None):
    rect(slide, 0, 0, W, Inches(1.3), NAVY)
    rect(slide, 0, Inches(1.3), W, Inches(0.06), BRIGHT_BLUE)
    txb(slide, title, Inches(0.5), Inches(0.18), Inches(12), Inches(0.85),
        size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if sub:
        txb(slide, sub, Inches(0.5), Inches(0.92), Inches(12), Inches(0.42),
            size=14, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

def screenshot_placeholder(slide, label="화면 캡처"):
    rect(slide, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
         line_color=MID_BLUE, line_width=1.5)
    txb(slide, f"[ {label} ]",
        Inches(6.75), Inches(4.0), Inches(6.1), Inches(1.0),
        size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ── Presentation ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, NAVY)

rect(s, 0, H - Inches(0.5),  W, Inches(0.5),  MID_BLUE)
rect(s, 0, H - Inches(0.55), W, Inches(0.08), BRIGHT_BLUE)

rect(s, Inches(0.55), Inches(0.55), Inches(2.7), Inches(0.44), BRIGHT_BLUE)
txb(s, "SparkLabs Korea", Inches(0.68), Inches(0.58),
    Inches(2.45), Inches(0.38), size=13, bold=True, color=WHITE)

txb(s, "서류 추적기",
    Inches(0.55), Inches(1.5), Inches(12), Inches(1.4),
    size=60, bold=True, color=WHITE)

txb(s, "AI 기반 투자 운용 자동화 플랫폼",
    Inches(0.55), Inches(2.9), Inches(12), Inches(0.75),
    size=24, color=LIGHT_BLUE)

rect(s, Inches(0.55), Inches(3.75), Inches(5), Inches(0.05), BRIGHT_BLUE)

txb(s, "2026년 8월  |  AI Agent 인턴십 최종 발표",
    Inches(0.55), Inches(3.95), Inches(10), Inches(0.55),
    size=15, color=LIGHT_BLUE)

rect(s, Inches(0.55), Inches(4.7), Inches(6.5), Inches(0.72), MID_BLUE)
txb(s, "11가지 기능   ·   Google Gemini API   ·   Next.js 16   ·   Vercel 배포",
    Inches(0.72), Inches(4.84), Inches(6.2), Inches(0.46),
    size=13, color=LIGHT_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "목차", "Table of Contents")

agenda = [
    ("01", "웹사이트 소개"),
    ("02", "기존 업무의 문제점"),
    ("03", "전체 기능 개요"),
    ("04", "주요 기능 상세  (6가지)"),
    ("05", "활용한 AI 도구"),
    ("06", "Gemini API 비용 분석"),
    ("07", "성과 및 개선 사항"),
    ("08", "한계점 · 향후 과제 · 라이브 데모"),
]
col_w = Inches(5.8)
for i, (num, label) in enumerate(agenda):
    col = i % 2
    row = i // 2
    lx = Inches(0.7) + col * Inches(6.6)
    ty = Inches(1.7) + row * Inches(1.3)
    rect(s, lx, ty, col_w, Inches(1.05), NAVY if col == 0 else MID_BLUE)
    txb(s, num, lx + Inches(0.18), ty + Inches(0.15), Inches(0.8), Inches(0.75),
        size=28, bold=True, color=BRIGHT_BLUE)
    txb(s, label, lx + Inches(0.72), ty + Inches(0.2),
        col_w - Inches(0.85), Inches(0.65), size=17, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 웹사이트 소개
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "웹사이트 소개", "SparkLabs 서류 추적기란?")

rect(s, Inches(0.5), Inches(1.55), Inches(7.6), Inches(5.5), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "SparkLabs 투자팀을 위한\nAll-in-One 내부 웹 플랫폼",
    Inches(0.75), Inches(1.7), Inches(7.1), Inches(0.95),
    size=20, bold=True, color=NAVY)
bullet_frame(s, [
    "투자 파이프라인 전 과정 자동화  (서류 수집 → 실사 → 계약 → 집행 → 전환)",
    "RCPS / SAFE 계약서 자동 생성  (원클릭 .docx 다운로드)",
    "Google Drive 연동 — 딜별 폴더 자동 생성 및 @sparklabs.co.kr 팀 공유",
    "Google Gemini AI — 파일 분류 · 실사 분석 · 챗봇 · 집행 추출 · 전환 계산",
    "한/영 이중 언어 지원  |  Google OAuth 로그인  |  다중 펀드 지원",
    "Vercel 배포 — 어디서든 접속 가능",
], Inches(0.75), Inches(2.75), Inches(7.1), Inches(4.0), size=15.5, color=DARK)

rect(s, Inches(8.4), Inches(1.55), Inches(4.45), Inches(2.75), NAVY)
txb(s, "Tech Stack", Inches(8.6), Inches(1.7), Inches(4.1), Inches(0.5),
    size=15, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "Next.js 16  ·  React 19  ·  TypeScript",
    "Tailwind CSS v4",
    "Vercel  ·  Vercel Blob (파일 저장)",
    "Google Drive API  ·  Google Gemini API",
    "NextAuth (Google OAuth 로그인)",
], Inches(8.6), Inches(2.28), Inches(4.1), Inches(1.9), size=13, color=WHITE,
   bullet_color=BRIGHT_BLUE)

rect(s, Inches(8.4), Inches(4.45), Inches(4.45), Inches(2.6), LIGHT_BLUE)
txb(s, "배포 URL", Inches(8.6), Inches(4.6), Inches(4.1), Inches(0.42),
    size=13, bold=True, color=NAVY)
txb(s, "sparklabs-doc-tracker.vercel.app",
    Inches(8.6), Inches(5.08), Inches(4.1), Inches(0.42),
    size=12, color=MID_BLUE, italic=True)
txb(s, "Google 계정 로그인 필수\n(@sparklabs.co.kr 내부 전용)",
    Inches(8.6), Inches(5.6), Inches(4.1), Inches(0.65),
    size=11.5, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 기존 업무의 문제점
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기존 업무의 문제점", "Before — What Made Investment Operations Painful?")

rect(s, Inches(0.5), Inches(1.55), Inches(3.6), Inches(5.5), NAVY)
txb(s, "계약서 1건당\n소요 시간",
    Inches(0.65), Inches(1.75), Inches(3.3), Inches(0.9),
    size=15, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
txb(s, "2~3\n시간",
    Inches(0.65), Inches(2.55), Inches(3.3), Inches(2.1),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "담당자 수동 작업 기준",
    Inches(0.65), Inches(4.75), Inches(3.3), Inches(0.5),
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(4.4), Inches(1.55), Inches(8.45), Inches(5.5), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "개발 전 주요 문제점",
    Inches(4.65), Inches(1.72), Inches(8.0), Inches(0.48),
    size=16, bold=True, color=NAVY)
bullet_frame(s, [
    "계약서 빈칸을 수동으로 하나하나 입력 — 동일한 정보를 수십 번 반복",
    "서류가 이메일 · 폴더 · Drive에 분산 — 중앙 추적 시스템 없음",
    "구조화된 실사 체크리스트 없음 — 항목별 완료 여부 추적 불가",
    "은행 30일 결제 마감 — 스프레드시트로 수동 관리, 누락 위험",
    "투자 집행 후 후속 서류 수집 프로세스가 비체계적",
    "업무 중 프로세스 질문 → 멘토 직접 문의 필요 (즉시 답변 불가)",
    "팀원마다 템플릿 버전 달라 일관성 부재",
], Inches(4.65), Inches(2.3), Inches(7.9), Inches(4.6),
   size=15.5, color=DARK, bullet_color=RED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 전체 기능 개요 (Pipeline Map)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "전체 기능 개요", "All Features — Investment Pipeline Map")

# Top global bar: Dashboard + Action Center
rect(s, Inches(0.5), Inches(1.44), Inches(12.35), Inches(0.6), BRIGHT_BLUE)
txb(s, "기능 ①  🏠 메인 대시보드  &  📋 액션 센터  —  전체 딜 현황 실시간 조회  ·  긴급 항목 우선순위 안내",
    Inches(0.65), Inches(1.52), Inches(12.0), Inches(0.42),
    size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Pipeline boxes
box_w  = Inches(1.65)
box_h  = Inches(1.6)
step   = Inches(2.02)   # box_w + 0.37 gap
sx     = Inches(0.84)
by     = Inches(2.18)

stages = [
    ("서류 수집",  MID_BLUE, "기능 ②", "드래그&드롭\nAI 분류\nDrive 연동"),
    ("서류 실사",  NAVY,     "기능 ③", "13항목\nAI 분석\n코멘트"),
    ("계약서 작성", MID_BLUE, "기능 ④", "RCPS 69토큰\nSAFE 40토큰\n.docx 다운"),
    ("투자 집행",  NAVY,     "기능 ⑤", "운용지시\n30일 마감\n이메일 초안"),
    ("SAFE 전환",  MID_BLUE, "기능 ⑥", "전환 추적\nAI 계산\n마감일 표시"),
    ("완료",       GREEN,    "",        ""),
]

for i, (name, col, num, sub) in enumerate(stages):
    lx = sx + i * step
    rect(s, lx, by, box_w, box_h, col)
    if num:
        txb(s, num, lx + Inches(0.07), by + Inches(0.06),
            box_w - Inches(0.1), Inches(0.26),
            size=9.5, bold=True, color=BRIGHT_BLUE if col == NAVY else LIGHT_BLUE)
    txb(s, name, lx + Inches(0.07), by + Inches(0.32),
        box_w - Inches(0.1), Inches(0.6),
        size=14.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        txb(s, sub, lx + Inches(0.07), by + Inches(0.95),
            box_w - Inches(0.1), Inches(0.65),
            size=9, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    # Arrow to next box
    if i < 5:
        arr_x = lx + box_w + Inches(0.04)
        arr_y = by + box_h / 2 - Inches(0.22)
        txb(s, "→", arr_x, arr_y, Inches(0.3), Inches(0.44),
            size=17, bold=True, color=BRIGHT_BLUE, align=PP_ALIGN.CENTER)

# Deal overview callout row
rect(s, Inches(0.5), Inches(3.94), Inches(12.35), Inches(0.58), DARK)
txb(s, "🌐  딜 개요 페이지  —  전체 6단계 진행 상황을 한 화면에서 조회  ·  단계별 링크 바로가기",
    Inches(0.65), Inches(4.04), Inches(12.0), Inches(0.38),
    size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# AI chatbot + bilingual bar
rect(s, Inches(0.5), Inches(4.62), Inches(12.35), Inches(0.58), MID_BLUE)
txb(s, "🤖  AI 프로세스 어시스턴트 (모든 페이지 우하단 플로팅)   ·   🌐  한/영 이중 언어 전체 지원   ·   🔐  Google OAuth",
    Inches(0.65), Inches(4.72), Inches(12.0), Inches(0.38),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Summary strip
rect(s, Inches(0.5), Inches(5.3), Inches(12.35), Inches(1.82), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "총 11가지 기능:",
    Inches(0.7), Inches(5.42), Inches(2.2), Inches(0.38),
    size=13, bold=True, color=NAVY)
txb(s, "대시보드  ·  액션 센터  ·  딜 개요  ·  서류 수집  ·  실사 체크리스트  ·  RCPS 계약서  ·  SAFE 계약서  ·  투자 집행  ·  SAFE 전환  ·  AI 챗봇  ·  이중 언어",
    Inches(3.1), Inches(5.42), Inches(9.6), Inches(0.38),
    size=13, color=DARK)
txb(s, "Google Gemini API  (gemini-flash-lite-latest)  —  파일 분류  ·  실사 분석  ·  계약서 필드 제안  ·  집행 번호 추출  ·  전환 계산  ·  챗봇 어시스턴트",
    Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.38),
    size=12, color=GRAY, italic=True)
txb(s, "Tech: Next.js 16  ·  React 19  ·  TypeScript  ·  Tailwind CSS v4  ·  Vercel  ·  Vercel Blob  ·  Google Drive API  ·  NextAuth",
    Inches(0.7), Inches(6.35), Inches(12.0), Inches(0.38),
    size=12, color=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 기능 ① 메인 대시보드 & 액션 센터
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기능 ①  메인 대시보드 & 액션 센터", "Main Dashboard & Action Center")

# Dashboard card
rect(s, Inches(0.5), Inches(1.55), Inches(5.9), Inches(2.6), NAVY)
txb(s, "🏠  메인 대시보드",
    Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.55),
    size=18, bold=True, color=WHITE)
bullet_frame(s, [
    "4개 KPI 타일:  추적 중인 회사 수  /  서류 준비 완료  /  DD 완료  /  주의 필요",
    "회사별 서류 수집 · 실사 진행률 바 (progress bar)",
    "딜 클릭 시 딜 개요 페이지로 바로 이동",
    "딜 개요:  6단계 전체 진행 상황 한 화면 요약",
], Inches(0.7), Inches(2.33), Inches(5.5), Inches(1.7),
   size=13.5, color=WHITE, bullet_color=BRIGHT_BLUE)

# Action Center card
rect(s, Inches(0.5), Inches(4.28), Inches(5.9), Inches(2.77), MID_BLUE)
txb(s, "📋  액션 센터 (Action Center)",
    Inches(0.7), Inches(4.43), Inches(5.5), Inches(0.55),
    size=18, bold=True, color=WHITE)
txb(s, '"지금 당장 뭘 해야 하나?"  우선순위 패널',
    Inches(0.7), Inches(5.0), Inches(5.5), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "① 은행 30일 결제 마감  (하드 컷오프, 날짜 강조 표시)",
    "② 필수 서류 미제출 항목  (딜명 + 미제출 서류 목록)",
    "③ 미완료 실사 체크 항목  (딜명 + 미완료 항목 수)",
], Inches(0.7), Inches(5.45), Inches(5.5), Inches(1.45),
   size=13.5, color=WHITE, bullet_color=ORANGE)

# Screenshot placeholder
rect(s, Inches(6.62), Inches(1.55), Inches(6.25), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ 화면 캡처 ]\n메인 대시보드",
    Inches(6.62), Inches(3.95), Inches(6.25), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 기능 ② 서류 수집 & AI 분류 & Drive 연동
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기능 ②  서류 수집 & AI 분류 & Drive 연동",
           "Document Collection, AI Classification & Google Drive")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), NAVY)
txb(s, "📂  서류 수집 탭",
    Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.55),
    size=18, bold=True, color=WHITE)
bullet_frame(s, [
    "드래그&드롭 파일 업로드 (drag-and-drop)",
    "딜별 필수 서류 체크리스트 — 업로드 시 자동 체크",
    "5초마다 자동 새로고침 (auto-poll, 실시간 반영)",
    "Vercel Blob에 비공개 안전 저장",
], Inches(0.7), Inches(2.33), Inches(5.6), Inches(1.55),
   size=14, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(3.95), Inches(5.4), Inches(0.06), BRIGHT_BLUE)
txb(s, "🤖  AI 파일명 자동 분류  (Google Gemini API)",
    Inches(0.7), Inches(4.08), Inches(5.6), Inches(0.42),
    size=13.5, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "키워드 매처가 인식 못 하는 파일명 → Gemini가 문서 종류 판별",
    "예: '최종본_v3.pdf'  →  등기부등본으로 자동 분류",
    "비용: ~$0.0002 / 호출  (연간 ~150 calls → ~$0.03)",
], Inches(0.7), Inches(4.57), Inches(5.6), Inches(0.98),
   size=13, color=LIGHT_BLUE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(5.62), Inches(5.4), Inches(0.06), BRIGHT_BLUE)
txb(s, "🗂️  Google Drive 연동",
    Inches(0.7), Inches(5.75), Inches(5.6), Inches(0.4),
    size=13.5, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "버튼 1회 클릭 → 딜별 폴더 자동 생성 + @sparklabs.co.kr 팀 공유",
    "업로드한 파일은 Drive 폴더에 자동 미러링",
], Inches(0.7), Inches(6.22), Inches(5.6), Inches(0.65),
   size=13, color=LIGHT_BLUE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ 화면 캡처 ]\n서류 수집 탭",
    Inches(6.75), Inches(3.95), Inches(6.1), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 기능 ③ 서류 실사 체크리스트
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기능 ③  서류 실사 체크리스트", "Due Diligence Checklist — 13 Points")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), MID_BLUE)
txb(s, "📋  서류 실사 탭",
    Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.55),
    size=18, bold=True, color=WHITE)
txb(s, "멘토 내부 문서 기반  13항목 체크리스트",
    Inches(0.7), Inches(2.28), Inches(5.6), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)

# Split: 1-7 document verification
rect(s, Inches(0.7), Inches(2.77), Inches(2.45), Inches(0.38), NAVY)
txb(s, "서류 검증  (1~7번)",
    Inches(0.8), Inches(2.8), Inches(2.25), Inches(0.32),
    size=11, bold=True, color=LIGHT_BLUE)
bullet_frame(s, [
    "등기부등본 / 사업자등록증",
    "주주명부 / 주식명세서",
    "재무제표 / 은행 잔고증명",
    "기타 필수 서류 확인",
], Inches(0.7), Inches(3.22), Inches(2.6), Inches(1.35),
   size=12, color=WHITE, bullet_color=BRIGHT_BLUE)

# Split: 8-13 process steps
rect(s, Inches(3.45), Inches(2.77), Inches(2.65), Inches(0.38), NAVY)
txb(s, "프로세스 절차  (8~13번)",
    Inches(3.55), Inches(2.8), Inches(2.45), Inches(0.32),
    size=11, bold=True, color=LIGHT_BLUE)
bullet_frame(s, [
    "주금 납입 확인",
    "등기 변경 여부 확인",
    "후속 조치 절차",
    "최종 완료 확인 사항",
], Inches(3.45), Inches(3.22), Inches(2.65), Inches(1.35),
   size=12, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(4.65), Inches(5.4), Inches(0.06), BRIGHT_BLUE)
txb(s, "추가 기능",
    Inches(0.7), Inches(4.78), Inches(5.6), Inches(0.38),
    size=13, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "항목별 메모 · 코멘트 직접 입력",
    "항목별 AI 분석 (Gemini) — 검토 포인트 제안",
    "완료 항목 자동 체크 및 전체 진행률 표시",
    "완료된 실사 현황이 대시보드 KPI에 반영",
], Inches(0.7), Inches(5.23), Inches(5.4), Inches(1.65),
   size=13, color=WHITE, bullet_color=LIGHT_BLUE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ 화면 캡처 ]\n서류 실사 체크리스트",
    Inches(6.75), Inches(3.95), Inches(6.1), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 기능 ④ 투자 계약서 자동 생성
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기능 ④  투자 계약서 자동 생성", "Investment Agreement Generator — RCPS & SAFE")

# RCPS panel
rect(s, Inches(0.5), Inches(1.55), Inches(3.6), Inches(5.5), NAVY)
txb(s, "📄  RCPS 계약서",
    Inches(0.65), Inches(1.7), Inches(3.35), Inches(0.55),
    size=16, bold=True, color=WHITE)
txb(s, "상환전환우선주 투자계약서",
    Inches(0.65), Inches(2.28), Inches(3.35), Inches(0.38),
    size=11.5, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "69개 빈칸 자동 채우기",
    "5개 섹션 구성\n(날짜/당사자/조건/통지/표준)",
    "서명란 · 별지1 · 별지2 포함",
    "표준 조항 기본값 내장\n(배당률, 상환이율 등)",
    "원클릭 .docx 다운로드",
], Inches(0.65), Inches(2.75), Inches(3.35), Inches(3.3),
   size=13, color=WHITE, bullet_color=BRIGHT_BLUE)
rect(s, Inches(0.65), Inches(6.3), Inches(3.3), Inches(0.5), BRIGHT_BLUE)
txb(s, "69 토큰 → 전체 자동 완성",
    Inches(0.7), Inches(6.36), Inches(3.2), Inches(0.38),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# SAFE panel
rect(s, Inches(4.28), Inches(1.55), Inches(3.6), Inches(5.5), MID_BLUE)
txb(s, "📄  SAFE 계약서",
    Inches(4.43), Inches(1.7), Inches(3.35), Inches(0.55),
    size=16, bold=True, color=WHITE)
txb(s, "조건부지분인수계약서",
    Inches(4.43), Inches(2.28), Inches(3.35), Inches(0.38),
    size=11.5, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "40개 빈칸 자동 채우기",
    "8개 섹션 구성",
    "Valuation Cap · 할인율",
    "별지1 · 별지3 동시 반영",
    "표준 조항 기본값 내장\n(위약벌 12%, 손배 120%)",
    "원클릭 .docx 다운로드",
], Inches(4.43), Inches(2.75), Inches(3.35), Inches(3.3),
   size=13, color=WHITE, bullet_color=LIGHT_BLUE)
rect(s, Inches(4.43), Inches(6.3), Inches(3.3), Inches(0.5), BRIGHT_BLUE)
txb(s, "40 토큰 → 전체 자동 완성",
    Inches(4.48), Inches(6.36), Inches(3.2), Inches(0.38),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Screenshot placeholder
rect(s, Inches(8.08), Inches(1.55), Inches(4.77), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ 화면 캡처 ]\n계약서 작성 탭\n(RCPS / SAFE)",
    Inches(8.08), Inches(3.7), Inches(4.77), Inches(1.4),
    size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 기능 ⑤ 투자 집행 추적
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기능 ⑤  투자 집행 추적", "Investment Execution Tracker")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), NAVY)
txb(s, "💼  투자 집행 탭",
    Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.55),
    size=18, bold=True, color=WHITE)
txb(s, "투자 집행 후 필수 절차를 단계별로 추적",
    Inches(0.7), Inches(2.28), Inches(5.6), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "운용지시 (bank payment instruction) 발송 및 추적",
    "주금 납입 후 후속 서류 수집 진행 관리",
    "수탁 은행 30일 결제 마감 — 하드 컷오프, 강조 표시",
    "펀드 구분:  모태펀드  vs  일반 사모펀드",
    "투자 구조 (투자 방식 · 금액 · 지분율) 기록",
    "이메일 초안 자동 생성 (AI 보조, 담당자 발송용)",
    "자동 저장 — 페이지 이탈 시 데이터 보존",
], Inches(0.7), Inches(2.75), Inches(5.6), Inches(3.7),
   size=14, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(6.55), Inches(5.4), Inches(0.32), RED)
txb(s, "⚠  은행 30일 마감 — 액션 센터에서 자동으로 최상위 노출",
    Inches(0.82), Inches(6.58), Inches(5.1), Inches(0.26),
    size=12, bold=True, color=WHITE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ 화면 캡처 ]\n투자 집행 추적 탭",
    Inches(6.75), Inches(3.95), Inches(6.1), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 기능 ⑥ SAFE 전환 추적 + AI 챗봇 어시스턴트
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "기능 ⑥  SAFE 전환 추적  +  AI 챗봇 어시스턴트",
           "SAFE Conversion Tracker & AI Process Assistant")

# SAFE Conversion (left)
rect(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.55), NAVY)
txb(s, "🔄  SAFE 전환 추적 탭",
    Inches(0.7), Inches(1.7), Inches(5.7), Inches(0.55),
    size=17, bold=True, color=WHITE)
txb(s, "SAFE → 주식 전환 프로세스 단계별 추적",
    Inches(0.7), Inches(2.28), Inches(5.7), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "전환 전 (Pre) · 전환 중 (During) · 전환 후 (Post) 3단계 구분",
    "마감일 자동 계산 및 D-day 표시",
    "AI 보조 전환 주식 수 계산 (Gemini)\n  — 투자 원금 · Valuation Cap · 할인율 기반",
    "전환 완료 후 등기 변경 사항 추적",
    "전환 프로세스 체크리스트 단계별 완료 표시",
], Inches(0.7), Inches(2.75), Inches(5.7), Inches(3.65),
   size=13.5, color=WHITE, bullet_color=BRIGHT_BLUE)

# AI Chatbot (right)
rect(s, Inches(6.85), Inches(1.55), Inches(6.0), Inches(5.55), MID_BLUE)
txb(s, "🤖  AI 프로세스 어시스턴트",
    Inches(7.05), Inches(1.7), Inches(5.6), Inches(0.55),
    size=17, bold=True, color=WHITE)
txb(s, "모든 페이지 우하단 — 플로팅 챗봇",
    Inches(7.05), Inches(2.28), Inches(5.6), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "스파크랩 내부 프로세스 문서 기반으로 학습",
    "현재 조회 중인 딜 정보 맥락 인식 (context-aware)",
    "한국어 · 영어 모두 답변 가능",
    "계약서 필드 입력 AI 제안 (Gemini)",
    "집행 관련 번호 자동 추출 (Gemini)",
    "멘토에게 묻지 않고 업무 중 즉시 가이드",
], Inches(7.05), Inches(2.75), Inches(5.6), Inches(3.65),
   size=13.5, color=WHITE, bullet_color=LIGHT_BLUE)

# Bilingual footer
rect(s, Inches(0.5), Inches(7.13), Inches(12.35), Inches(0.35), DARK)
txb(s, "🌐  한/영 이중 언어 — 플랫폼 전체 언어 토글 지원  |  UI 전체 한국어/영어 전환 가능",
    Inches(0.65), Inches(7.17), Inches(12.0), Inches(0.28),
    size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — 활용한 AI 도구
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "활용한 AI 도구", "AI Tools — Two Different Roles")

# Left: Gemini in the app
rect(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.55), NAVY)
txb(s, "앱 내부에서 사용한 AI",
    Inches(0.7), Inches(1.72), Inches(5.7), Inches(0.38),
    size=12, bold=True, color=BRIGHT_BLUE)
txb(s, "Google Gemini API",
    Inches(0.7), Inches(2.13), Inches(5.7), Inches(0.6),
    size=22, bold=True, color=WHITE)
txb(s, "gemini-flash-lite-latest",
    Inches(0.7), Inches(2.75), Inches(5.7), Inches(0.38),
    size=12, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "파일명 자동 분류  (서류 수집 탭)",
    "서류 실사 AI 분석 — 항목별 검토 포인트 제안",
    "계약서 필드 AI 제안  (계약서 탭)",
    "집행 관련 번호 자동 추출  (투자 집행 탭)",
    "SAFE 전환 주식 수 계산 보조  (전환 추적 탭)",
    "인앱 챗봇 어시스턴트 — 프로세스 Q&A",
], Inches(0.7), Inches(3.2), Inches(5.6), Inches(3.55),
   size=13.5, color=WHITE, bullet_color=BRIGHT_BLUE)

# Right: Claude used to build
rect(s, Inches(6.85), Inches(1.55), Inches(6.0), Inches(5.55), MID_BLUE)
txb(s, "개발에 활용한 AI 도구",
    Inches(7.05), Inches(1.72), Inches(5.6), Inches(0.38),
    size=12, bold=True, color=LIGHT_BLUE)

rect(s, Inches(7.05), Inches(2.15), Inches(5.6), Inches(1.85), NAVY)
txb(s, "Claude  (Anthropic)",
    Inches(7.25), Inches(2.27), Inches(5.2), Inches(0.5),
    size=17, bold=True, color=WHITE)
txb(s, "코드 작성, 디버깅, 기능 설계\nDOCX XML 조작 등 복잡한 로직 구현",
    Inches(7.25), Inches(2.77), Inches(5.2), Inches(1.0),
    size=13, color=LIGHT_BLUE)

rect(s, Inches(7.05), Inches(4.13), Inches(5.6), Inches(1.85), NAVY)
txb(s, "Claude Code  (CLI 도구)",
    Inches(7.25), Inches(4.25), Inches(5.2), Inches(0.5),
    size=17, bold=True, color=WHITE)
txb(s, "터미널에서 파일 읽기 · 편집 · 실행\n프로그래밍 경험 없이 개발 가능하게 해줌",
    Inches(7.25), Inches(4.75), Inches(5.2), Inches(1.0),
    size=13, color=LIGHT_BLUE)

rect(s, 0, Inches(7.1), W, Inches(0.4), LIGHT_BLUE)
txb(s, "전체 스택:   Next.js 16  ·  React 19  ·  TypeScript  ·  Tailwind CSS v4  ·  Vercel  ·  Vercel Blob  ·  Google Drive API  ·  Google Gemini API  ·  NextAuth",
    Inches(0.7), Inches(7.13), Inches(12.5), Inches(0.35),
    size=12, color=NAVY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Google Gemini API 비용 분석
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Google Gemini API — 비용 분석", "How Much Does the AI Actually Cost?")

# Left: model info
rect(s, Inches(0.5), Inches(1.55), Inches(3.7), Inches(5.55), NAVY)
txb(s, "사용 모델",
    Inches(0.7), Inches(1.72), Inches(3.3), Inches(0.38),
    size=12, bold=True, color=BRIGHT_BLUE)
txb(s, "gemini-flash-\nlite-latest",
    Inches(0.7), Inches(2.12), Inches(3.3), Inches(1.1),
    size=20, bold=True, color=WHITE)
txb(s, "Google의 가장 저렴한\n최신 AI 모델",
    Inches(0.7), Inches(3.22), Inches(3.3), Inches(0.65),
    size=12, color=LIGHT_BLUE)
rect(s, Inches(0.7), Inches(4.0), Inches(3.1), Inches(0.06), BRIGHT_BLUE)
txb(s, "호출당 비용",
    Inches(0.7), Inches(4.15), Inches(3.3), Inches(0.38),
    size=12, bold=True, color=BRIGHT_BLUE)
txb(s, "~$0.0002",
    Inches(0.7), Inches(4.55), Inches(3.3), Inches(0.75),
    size=26, bold=True, color=WHITE)
txb(s, "per API call",
    Inches(0.7), Inches(5.3), Inches(3.3), Inches(0.38),
    size=12, color=LIGHT_BLUE, italic=True)
txb(s, "Claude Haiku보다 13배 저렴",
    Inches(0.7), Inches(5.78), Inches(3.3), Inches(0.38),
    size=11, color=GRAY)

# Centre: usage table
rect(s, Inches(4.45), Inches(1.55), Inches(5.3), Inches(5.55), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "연간 예상 사용량",
    Inches(4.65), Inches(1.72), Inches(4.9), Inches(0.42),
    size=14, bold=True, color=NAVY)

rows = [
    ("파일명 자동 분류",   "~150 calls", "~$0.03"),
    ("실사 AI 분석",       "~50 calls",  "~$0.01"),
    ("챗봇 어시스턴트",    "~250 calls", "~$0.05"),
]
row_cols = [LIGHT_BLUE, OFF_WHITE, LIGHT_BLUE]
for i, (label, calls, cost) in enumerate(rows):
    ty = Inches(2.28) + i * Inches(0.77)
    rect(s, Inches(4.55), ty, Inches(5.1), Inches(0.65), row_cols[i])
    txb(s, label, Inches(4.72), ty + Inches(0.12), Inches(2.4), Inches(0.4),
        size=13, color=DARK)
    txb(s, calls, Inches(7.2), ty + Inches(0.12), Inches(1.0), Inches(0.4),
        size=13, color=GRAY, align=PP_ALIGN.CENTER)
    txb(s, cost, Inches(8.25), ty + Inches(0.12), Inches(1.2), Inches(0.4),
        size=13, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

rect(s, Inches(4.55), Inches(4.59), Inches(5.1), Inches(0.82), NAVY)
txb(s, "연간 합계", Inches(4.72), Inches(4.72), Inches(2.4), Inches(0.5),
    size=14, bold=True, color=WHITE)
txb(s, "~450 calls", Inches(7.2), Inches(4.72), Inches(1.0), Inches(0.5),
    size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
txb(s, "~$0.09", Inches(8.05), Inches(4.69), Inches(1.4), Inches(0.55),
    size=20, bold=True, color=BRIGHT_BLUE, align=PP_ALIGN.RIGHT)

txb(s, "* 10배 사용:  ~$0.90/year   *  100배 사용:  ~$9.00/year",
    Inches(4.65), Inches(5.55), Inches(5.0), Inches(0.38),
    size=11, color=GRAY, italic=True)
txb(s, "무료 티어: 분당 60회 요청 무료\n(현재 사용량 기준 무료 범위 안에 해당)",
    Inches(4.65), Inches(6.0), Inches(5.0), Inches(0.65),
    size=12, color=DARK)

# Right: dramatic green callout
rect(s, Inches(10.0), Inches(1.55), Inches(2.85), Inches(5.55), GREEN)
txb(s, "연간 AI\n사용료",
    Inches(10.1), Inches(1.72), Inches(2.6), Inches(0.85),
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "$0.09",
    Inches(10.1), Inches(2.55), Inches(2.6), Inches(1.5),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "per year",
    Inches(10.1), Inches(4.05), Inches(2.6), Inches(0.45),
    size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
rect(s, Inches(10.2), Inches(4.6), Inches(2.4), Inches(0.05), WHITE)
txb(s, "스타벅스\n아메리카노",
    Inches(10.1), Inches(4.75), Inches(2.6), Inches(0.75),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "~$5.00",
    Inches(10.1), Inches(5.5), Inches(2.6), Inches(0.65),
    size=20, bold=True, color=rgb(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)
txb(s, "커피 한 잔 > AI 1년",
    Inches(10.1), Inches(6.2), Inches(2.6), Inches(0.45),
    size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — 성과 및 개선 사항
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, NAVY)
header_bar(s, "성과 및 개선 사항", "Results & Improvements")
rect(s, 0, Inches(1.3), W, Inches(0.06), BRIGHT_BLUE)

# Before
rect(s, Inches(0.5), Inches(1.6), Inches(5.5), Inches(3.75), RED)
txb(s, "이전 (Before)",
    Inches(0.7), Inches(1.75), Inches(5.1), Inches(0.52),
    size=17, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "2~3\n시간",
    Inches(0.7), Inches(2.28), Inches(5.1), Inches(1.5),
    size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "계약서 1건 작성 기준\n(수동 입력 + 검토)",
    Inches(0.7), Inches(3.82), Inches(5.1), Inches(0.7),
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "↓  담당자가 모든 빈칸 수동 입력",
    Inches(0.7), Inches(4.65), Inches(5.1), Inches(0.48),
    size=12, color=rgb(255, 200, 200), align=PP_ALIGN.CENTER)

# Arrow
txb(s, "→", Inches(5.5), Inches(2.9), Inches(0.9), Inches(1.5),
    size=50, bold=True, color=BRIGHT_BLUE, align=PP_ALIGN.CENTER)

# After
rect(s, Inches(6.3), Inches(1.6), Inches(6.55), Inches(3.75), GREEN)
txb(s, "현재 (After)",
    Inches(6.5), Inches(1.75), Inches(6.2), Inches(0.52),
    size=17, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "~3분",
    Inches(6.5), Inches(2.28), Inches(6.2), Inches(1.5),
    size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "계약서 1건 작성 기준\n(정보 입력 후 즉시 다운로드)",
    Inches(6.5), Inches(3.82), Inches(6.2), Inches(0.7),
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "↓  69개 빈칸 전체 자동 완성",
    Inches(6.5), Inches(4.65), Inches(6.2), Inches(0.48),
    size=12, color=rgb(200, 255, 200), align=PP_ALIGN.CENTER)

# Additional improvements bar
rect(s, Inches(0.5), Inches(5.52), Inches(12.35), Inches(0.48), MID_BLUE)
txb(s, "추가 성과:  은행 마감 자동 알림  ·  서류 중앙 추적  ·  실사 체계화  ·  AI 즉시 답변  ·  이중 언어 지원",
    Inches(0.65), Inches(5.6), Inches(12.0), Inches(0.35),
    size=13, color=WHITE, align=PP_ALIGN.CENTER)

# Big stat bar
rect(s, Inches(0.5), Inches(6.12), Inches(12.35), Inches(1.0), BRIGHT_BLUE)
txb(s, "약 97% 이상 시간 단축  (2~3시간 → ~3분)",
    Inches(0.7), Inches(6.18), Inches(12.0), Inches(0.82),
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — 한계점 및 향후 과제
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "한계점 및 향후 과제", "Current Limitations & Future Roadmap")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), NAVY)
txb(s, "⚠️  현재 한계점",
    Inches(0.7), Inches(1.72), Inches(5.5), Inches(0.58),
    size=20, bold=True, color=ORANGE)
bullet_frame(s, [
    "계약서 템플릿 업데이트 수동\n  (법적 조항 변경 시 XML 직접 수정)",
    "다수의 이해관계인 처리 미구현 (현재 1명만 지원)",
    "모바일 최적화 미완성",
    "전자서명 기능 없음 (현재 습식 서명 필요)",
    "사용자 권한 관리 없음 (팀원 전체 동일 권한)",
    "실사 체크리스트 항목 수 조정 어려움",
], Inches(0.7), Inches(2.45), Inches(5.5), Inches(4.4),
   size=14, color=WHITE, bullet_color=ORANGE)

rect(s, Inches(6.8), Inches(1.55), Inches(6.05), Inches(5.5), LIGHT_BLUE)
txb(s, "🚀  앞으로 추가할 기능",
    Inches(7.0), Inches(1.72), Inches(5.6), Inches(0.58),
    size=20, bold=True, color=NAVY)
bullet_frame(s, [
    "이메일 알림 (계약서 발송 · 서명 완료)",
    "전자서명 연동  (e.g. DocuSign)",
    "AI 기반 계약서 조항 자동 검토",
    "펀드별 투자 통계 대시보드",
    "다중 이해관계인 지원",
    "계약서 버전 관리 (개정 이력 추적)",
    "관리자용 템플릿 UI 편집기",
], Inches(7.0), Inches(2.45), Inches(5.6), Inches(4.4),
   size=14, color=DARK, bullet_color=MID_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — 라이브 데모
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, NAVY)

rect(s, 0, H - Inches(0.55), W, Inches(0.55), BRIGHT_BLUE)
rect(s, 0, Inches(0),        W, Inches(0.55), BRIGHT_BLUE)

txb(s, "라이브 데모",
    Inches(1.5), Inches(1.8), Inches(10), Inches(2.0),
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Live Demonstration",
    Inches(1.5), Inches(3.82), Inches(10), Inches(0.75),
    size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(4.78), Inches(6.35), Inches(0.08), BRIGHT_BLUE)

txb(s, "실제 사이트를 직접 시연합니다",
    Inches(1.5), Inches(5.02), Inches(10), Inches(0.6),
    size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

demo_items = ["대시보드 조회", "서류 수집 탭", "계약서 작성 & 다운로드", "AI 챗봇 시연"]
txb(s, "     ·     ".join(demo_items),
    Inches(0.5), Inches(5.88), Inches(12.35), Inches(0.55),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "sparklabs-presentation-v2.pptx")
prs.save(out)
print(f"Saved: {out}")
