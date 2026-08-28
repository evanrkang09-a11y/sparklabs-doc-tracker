"""
SparkLabs Document Tracker — English Presentation Generator (v2)
16 slides, 15-20 minutes
Output: sparklabs-presentation-english.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x2B, 0x55)
MID_BLUE    = RGBColor(0x1A, 0x56, 0xAA)
BRIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xFF)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
DARK        = RGBColor(0x1E, 0x29, 0x3B)
GRAY        = RGBColor(0x64, 0x74, 0x8B)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)

def rgb(r, g, b): return RGBColor(r, g, b)

W, H = Inches(13.33), Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line_color=None, line_width=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp

def txb(slide, text, l, t, w, h, size, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box

def bullet_frame(slide, items, l, t, w, h, size=17, color=DARK,
                 bullet_color=BRIGHT_BLUE, line_space=1.2, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = para._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_space * 100000)))
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", "100")
        br = para.add_run()
        br.text = "● "
        br.font.size = Pt(size * 0.62)
        br.font.color.rgb = bullet_color
        br.font.name = font
        tr = para.add_run()
        if isinstance(item, tuple):
            tr.text = item[0]
            tr.font.size = Pt(item[1] if len(item) > 1 else size)
            tr.font.bold = item[2] if len(item) > 2 else False
            tr.font.color.rgb = item[3] if len(item) > 3 else color
        else:
            tr.text = item
            tr.font.size = Pt(size)
            tr.font.color.rgb = color
        tr.font.name = font
    return box

def header_bar(slide, title, sub=None):
    rect(slide, 0, 0, W, Inches(1.3), NAVY)
    rect(slide, 0, Inches(1.3), W, Inches(0.06), BRIGHT_BLUE)
    txb(slide, title, Inches(0.5), Inches(0.18), Inches(12), Inches(0.85),
        size=30, bold=True, color=WHITE)
    if sub:
        txb(slide, sub, Inches(0.5), Inches(0.92), Inches(12), Inches(0.42),
            size=14, color=LIGHT_BLUE)

# ── Presentation ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, NAVY)

rect(s, 0, H - Inches(0.5),  W, Inches(0.5),  MID_BLUE)
rect(s, 0, H - Inches(0.55), W, Inches(0.08), BRIGHT_BLUE)

rect(s, Inches(0.55), Inches(0.55), Inches(2.7), Inches(0.44), BRIGHT_BLUE)
txb(s, "SparkLabs Korea", Inches(0.68), Inches(0.58),
    Inches(2.45), Inches(0.38), size=13, bold=True, color=WHITE)

txb(s, "Document Tracker",
    Inches(0.55), Inches(1.5), Inches(12), Inches(1.4),
    size=60, bold=True, color=WHITE)

txb(s, "AI-Powered Investment Operations Automation Platform",
    Inches(0.55), Inches(2.9), Inches(12), Inches(0.75),
    size=24, color=LIGHT_BLUE)

rect(s, Inches(0.55), Inches(3.75), Inches(5), Inches(0.05), BRIGHT_BLUE)

txb(s, "August 2026  |  AI Agent Internship — Final Presentation",
    Inches(0.55), Inches(3.95), Inches(10), Inches(0.55),
    size=15, color=LIGHT_BLUE)

rect(s, Inches(0.55), Inches(4.7), Inches(6.5), Inches(0.72), MID_BLUE)
txb(s, "11 features   ·   Google Gemini API   ·   Next.js 16   ·   Deployed on Vercel",
    Inches(0.72), Inches(4.84), Inches(6.2), Inches(0.46),
    size=13, color=LIGHT_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Agenda")

agenda = [
    ("01", "Website Overview"),
    ("02", "The Problem — Before"),
    ("03", "All Features Overview"),
    ("04", "Feature Deep Dives  (6 features)"),
    ("05", "AI Tools Used"),
    ("06", "Gemini API Cost Analysis"),
    ("07", "Results & Improvements"),
    ("08", "Limitations · Roadmap · Live Demo"),
]
col_w = Inches(5.8)
for i, (num, label) in enumerate(agenda):
    col = i % 2
    row = i // 2
    lx = Inches(0.7) + col * Inches(6.6)
    ty = Inches(1.7) + row * Inches(1.3)
    rect(s, lx, ty, col_w, Inches(1.05), NAVY if col == 0 else MID_BLUE)
    txb(s, num, lx + Inches(0.18), ty + Inches(0.15), Inches(0.8), Inches(0.75),
        size=28, bold=True, color=BRIGHT_BLUE)
    txb(s, label, lx + Inches(0.72), ty + Inches(0.2),
        col_w - Inches(0.85), Inches(0.65), size=17, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WEBSITE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Website Overview", "What is the SparkLabs Document Tracker?")

rect(s, Inches(0.5), Inches(1.55), Inches(7.6), Inches(5.5), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "An All-in-One internal web platform\nfor the SparkLabs investment team",
    Inches(0.75), Inches(1.7), Inches(7.1), Inches(0.95),
    size=20, bold=True, color=NAVY)
bullet_frame(s, [
    "Automates the full investment pipeline  (doc collection → diligence → agreement → execution → conversion)",
    "Auto-generates RCPS / SAFE agreements  (one-click .docx download)",
    "Google Drive integration — auto-creates a deal folder and shares with @sparklabs.co.kr",
    "Google Gemini AI — file classification · diligence analysis · chatbot · execution extraction · conversion calc",
    "Full KO/EN bilingual support  |  Google OAuth login  |  Multi-fund support",
    "Deployed on Vercel — accessible from anywhere",
], Inches(0.75), Inches(2.75), Inches(7.1), Inches(4.0), size=15.5, color=DARK)

rect(s, Inches(8.4), Inches(1.55), Inches(4.45), Inches(2.75), NAVY)
txb(s, "Tech Stack", Inches(8.6), Inches(1.7), Inches(4.1), Inches(0.5),
    size=15, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "Next.js 16  ·  React 19  ·  TypeScript",
    "Tailwind CSS v4",
    "Vercel  ·  Vercel Blob (file storage)",
    "Google Drive API  ·  Google Gemini API",
    "NextAuth (Google OAuth login)",
], Inches(8.6), Inches(2.28), Inches(4.1), Inches(1.9), size=13, color=WHITE,
   bullet_color=BRIGHT_BLUE)

rect(s, Inches(8.4), Inches(4.45), Inches(4.45), Inches(2.6), LIGHT_BLUE)
txb(s, "Deployed at", Inches(8.6), Inches(4.6), Inches(4.1), Inches(0.42),
    size=13, bold=True, color=NAVY)
txb(s, "sparklabs-doc-tracker.vercel.app",
    Inches(8.6), Inches(5.08), Inches(4.1), Inches(0.42),
    size=12, color=MID_BLUE, italic=True)
txb(s, "Google account login required\n(@sparklabs.co.kr — internal only)",
    Inches(8.6), Inches(5.6), Inches(4.1), Inches(0.65),
    size=11.5, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "The Problem — Before Development")

rect(s, Inches(0.5), Inches(1.55), Inches(3.6), Inches(5.5), NAVY)
txb(s, "Time per\nagreement",
    Inches(0.65), Inches(1.75), Inches(3.3), Inches(0.9),
    size=15, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
txb(s, "2~3\nhours",
    Inches(0.65), Inches(2.55), Inches(3.3), Inches(2.1),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "per team member, per deal",
    Inches(0.65), Inches(4.75), Inches(3.3), Inches(0.5),
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(4.4), Inches(1.55), Inches(8.45), Inches(5.5), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "Key problems before the platform was built",
    Inches(4.65), Inches(1.72), Inches(8.0), Inches(0.48),
    size=16, bold=True, color=NAVY)
bullet_frame(s, [
    "Fill every blank in the Word template by hand — same info repeated dozens of times",
    "Documents scattered across emails, folders, and Drive — no central tracking",
    "No structured due diligence checklist — impossible to track item-by-item completion",
    "30-day bank payment deadline tracked manually in spreadsheets — risk of missing it",
    "No structured process for collecting post-execution documents",
    "Process questions during work → had to ask mentor directly (no instant answer)",
    "Each team member's template version may differ — no consistency",
], Inches(4.65), Inches(2.3), Inches(7.9), Inches(4.6),
   size=15.5, color=DARK, bullet_color=RED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ALL FEATURES OVERVIEW (Pipeline Map)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "All Features Overview", "Investment Pipeline Map")

# Top global bar
rect(s, Inches(0.5), Inches(1.44), Inches(12.35), Inches(0.6), BRIGHT_BLUE)
txb(s, "Feature 1  —  Main Dashboard  &  Action Center  —  Real-time deal overview  ·  Priority urgent items",
    Inches(0.65), Inches(1.52), Inches(12.0), Inches(0.42),
    size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box_w  = Inches(1.65)
box_h  = Inches(1.6)
step   = Inches(2.02)
sx     = Inches(0.84)
by     = Inches(2.18)

stages = [
    ("Doc\nCollection",  MID_BLUE, "Feature 2", "Drag & drop\nAI classify\nDrive sync"),
    ("Due\nDiligence",   NAVY,     "Feature 3", "13-point list\nAI analysis\nComments"),
    ("Agreement\nDrafting", MID_BLUE, "Feature 4", "RCPS 69 tokens\nSAFE 40 tokens\n.docx download"),
    ("Execution\nTracking",  NAVY,   "Feature 5", "Bank instruction\n30-day deadline\nEmail draft"),
    ("SAFE\nConversion",  MID_BLUE, "Feature 6", "Pre/During/Post\nAI calculation\nDeadline calc"),
    ("Complete", GREEN,    "",        ""),
]

for i, (name, col, num, sub) in enumerate(stages):
    lx = sx + i * step
    rect(s, lx, by, box_w, box_h, col)
    if num:
        txb(s, num, lx + Inches(0.07), by + Inches(0.06),
            box_w - Inches(0.1), Inches(0.26),
            size=9.5, bold=True, color=BRIGHT_BLUE if col == NAVY else LIGHT_BLUE)
    txb(s, name, lx + Inches(0.07), by + Inches(0.32),
        box_w - Inches(0.1), Inches(0.6),
        size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        txb(s, sub, lx + Inches(0.07), by + Inches(0.95),
            box_w - Inches(0.1), Inches(0.65),
            size=9, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    if i < 5:
        arr_x = lx + box_w + Inches(0.04)
        arr_y = by + box_h / 2 - Inches(0.22)
        txb(s, "→", arr_x, arr_y, Inches(0.3), Inches(0.44),
            size=17, bold=True, color=BRIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(3.94), Inches(12.35), Inches(0.58), DARK)
txb(s, "Deal Overview page — full 6-stage status summary on one screen  ·  direct links to each stage",
    Inches(0.65), Inches(4.04), Inches(12.0), Inches(0.38),
    size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(4.62), Inches(12.35), Inches(0.58), MID_BLUE)
txb(s, "AI Process Assistant (floating chatbot, every page)   ·   Full KO/EN bilingual toggle   ·   Google OAuth login",
    Inches(0.65), Inches(4.72), Inches(12.0), Inches(0.38),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(5.3), Inches(12.35), Inches(1.82), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "11 total features:",
    Inches(0.7), Inches(5.42), Inches(2.2), Inches(0.38),
    size=13, bold=True, color=NAVY)
txb(s, "Dashboard  ·  Action Center  ·  Deal Overview  ·  Doc Collection  ·  Diligence Checklist  ·  RCPS Agreement  ·  SAFE Agreement  ·  Exec Tracker  ·  SAFE Conversion  ·  AI Chatbot  ·  Bilingual",
    Inches(3.1), Inches(5.42), Inches(9.6), Inches(0.38),
    size=13, color=DARK)
txb(s, "Google Gemini API  (gemini-flash-lite-latest)  —  file classification  ·  diligence analysis  ·  agreement suggestions  ·  number extraction  ·  conversion calc  ·  chatbot",
    Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.38),
    size=12, color=GRAY, italic=True)
txb(s, "Tech: Next.js 16  ·  React 19  ·  TypeScript  ·  Tailwind CSS v4  ·  Vercel  ·  Vercel Blob  ·  Google Drive API  ·  NextAuth",
    Inches(0.7), Inches(6.35), Inches(12.0), Inches(0.38),
    size=12, color=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURE 1: MAIN DASHBOARD & ACTION CENTER
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Feature 1  —  Main Dashboard & Action Center")

rect(s, Inches(0.5), Inches(1.55), Inches(5.9), Inches(2.6), NAVY)
txb(s, "Main Dashboard",
    Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.55),
    size=18, bold=True, color=WHITE)
bullet_frame(s, [
    "4 KPI tiles:  Companies tracked  /  Docs ready  /  DD complete  /  Needs attention",
    "Per-company progress bars for document collection and diligence",
    "Click a deal → go directly to the Deal Overview page",
    "Deal Overview:  one-screen summary of all 6 pipeline stages",
], Inches(0.7), Inches(2.33), Inches(5.5), Inches(1.7),
   size=13.5, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.5), Inches(4.28), Inches(5.9), Inches(2.77), MID_BLUE)
txb(s, "Action Center",
    Inches(0.7), Inches(4.43), Inches(5.5), Inches(0.55),
    size=18, bold=True, color=WHITE)
txb(s, '"What do I need to do right now?"  — Priority panel',
    Inches(0.7), Inches(5.0), Inches(5.5), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "Priority 1 — 30-day bank payment deadline  (hard cutoff, date highlighted)",
    "Priority 2 — Missing required documents  (deal name + missing docs list)",
    "Priority 3 — Outstanding diligence items  (deal name + item count)",
], Inches(0.7), Inches(5.45), Inches(5.5), Inches(1.45),
   size=13.5, color=WHITE, bullet_color=ORANGE)

rect(s, Inches(6.62), Inches(1.55), Inches(6.25), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ Screenshot ]\nMain Dashboard",
    Inches(6.62), Inches(3.95), Inches(6.25), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURE 2: DOCUMENT COLLECTION, AI CLASSIFICATION & DRIVE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Feature 2  —  Document Collection, AI Classification & Drive",
           "서류 수집 Tab")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), NAVY)
txb(s, "Document Collection Tab",
    Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.55),
    size=18, bold=True, color=WHITE)
bullet_frame(s, [
    "Drag-and-drop file upload",
    "Required document checklist per deal — auto-checks when file is uploaded",
    "Auto-refreshes every 5 seconds (real-time polling)",
    "Files stored privately on Vercel Blob",
], Inches(0.7), Inches(2.33), Inches(5.6), Inches(1.55),
   size=14, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(3.95), Inches(5.4), Inches(0.06), BRIGHT_BLUE)
txb(s, "AI Filename Auto-Classification  (Google Gemini API)",
    Inches(0.7), Inches(4.08), Inches(5.6), Inches(0.42),
    size=13.5, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "When keyword matching fails, Gemini identifies the document type",
    "Example: 'final_v3.pdf'  →  auto-classified as certificate of registration",
    "Cost: ~$0.0002 / call  (~150 calls/year → ~$0.03/year)",
], Inches(0.7), Inches(4.57), Inches(5.6), Inches(0.98),
   size=13, color=LIGHT_BLUE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(5.62), Inches(5.4), Inches(0.06), BRIGHT_BLUE)
txb(s, "Google Drive Integration",
    Inches(0.7), Inches(5.75), Inches(5.6), Inches(0.4),
    size=13.5, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "One button click → deal folder auto-created and shared with @sparklabs.co.kr team",
    "Uploaded files automatically mirrored to the Drive folder",
], Inches(0.7), Inches(6.22), Inches(5.6), Inches(0.65),
   size=13, color=LIGHT_BLUE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ Screenshot ]\nDocument Collection Tab",
    Inches(6.75), Inches(3.95), Inches(6.1), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURE 3: DUE DILIGENCE CHECKLIST
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Feature 3  —  Due Diligence Checklist", "서류 실사 Tab — 13-Point Checklist")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), MID_BLUE)
txb(s, "Due Diligence Checklist Tab",
    Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.55),
    size=18, bold=True, color=WHITE)
txb(s, "13-point checklist based on mentor's internal process document",
    Inches(0.7), Inches(2.28), Inches(5.6), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)

rect(s, Inches(0.7), Inches(2.77), Inches(2.45), Inches(0.38), NAVY)
txb(s, "Document Verification  (1–7)",
    Inches(0.8), Inches(2.8), Inches(2.25), Inches(0.32),
    size=11, bold=True, color=LIGHT_BLUE)
bullet_frame(s, [
    "Certificate of registration",
    "Business registration",
    "Shareholder register / share register",
    "Financial statements / bank balance",
], Inches(0.7), Inches(3.22), Inches(2.6), Inches(1.35),
   size=12, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(3.45), Inches(2.77), Inches(2.65), Inches(0.38), NAVY)
txb(s, "Process Steps  (8–13)",
    Inches(3.55), Inches(2.8), Inches(2.45), Inches(0.32),
    size=11, bold=True, color=LIGHT_BLUE)
bullet_frame(s, [
    "Confirm payment deposit",
    "Verify registry changes",
    "Follow-up procedures",
    "Final completion checks",
], Inches(3.45), Inches(3.22), Inches(2.65), Inches(1.35),
   size=12, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(4.65), Inches(5.4), Inches(0.06), BRIGHT_BLUE)
txb(s, "Additional Features",
    Inches(0.7), Inches(4.78), Inches(5.6), Inches(0.38),
    size=13, bold=True, color=BRIGHT_BLUE)
bullet_frame(s, [
    "Per-item notes and comments",
    "Per-item AI analysis (Gemini) — review point suggestions",
    "Auto-check completed items and display overall progress",
    "Completion status feeds back into the dashboard KPI tiles",
], Inches(0.7), Inches(5.23), Inches(5.4), Inches(1.65),
   size=13, color=WHITE, bullet_color=LIGHT_BLUE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ Screenshot ]\nDue Diligence Checklist",
    Inches(6.75), Inches(3.95), Inches(6.1), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FEATURE 4: INVESTMENT AGREEMENT GENERATOR
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Feature 4  —  Investment Agreement Generator", "RCPS & SAFE Auto-Generation")

# RCPS
rect(s, Inches(0.5), Inches(1.55), Inches(3.6), Inches(5.5), NAVY)
txb(s, "RCPS Agreement",
    Inches(0.65), Inches(1.7), Inches(3.35), Inches(0.55),
    size=16, bold=True, color=WHITE)
txb(s, "Redeemable Convertible Preferred Stock",
    Inches(0.65), Inches(2.28), Inches(3.35), Inches(0.38),
    size=11, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "69 blank fields — all filled automatically",
    "5 sections\n(Dates / Parties / Terms / Notices / Standard)",
    "Signature block + Appendix 1 & 2",
    "SparkLabs standard defaults pre-filled\n(dividend rate, redemption rate…)",
    "One-click .docx download",
], Inches(0.65), Inches(2.75), Inches(3.35), Inches(3.3),
   size=13, color=WHITE, bullet_color=BRIGHT_BLUE)
rect(s, Inches(0.65), Inches(6.3), Inches(3.3), Inches(0.5), BRIGHT_BLUE)
txb(s, "69 tokens → auto-completed",
    Inches(0.7), Inches(6.36), Inches(3.2), Inches(0.38),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# SAFE
rect(s, Inches(4.28), Inches(1.55), Inches(3.6), Inches(5.5), MID_BLUE)
txb(s, "SAFE Agreement",
    Inches(4.43), Inches(1.7), Inches(3.35), Inches(0.55),
    size=16, bold=True, color=WHITE)
txb(s, "Simple Agreement for Future Equity",
    Inches(4.43), Inches(2.28), Inches(3.35), Inches(0.38),
    size=11, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "40 blank fields — all filled automatically",
    "8 sections",
    "Valuation cap · discount rate",
    "Interested party fills Appendix 1 & 3 simultaneously",
    "Standard clauses pre-filled\n(12% penalty, 120% damages…)",
    "One-click .docx download",
], Inches(4.43), Inches(2.75), Inches(3.35), Inches(3.3),
   size=13, color=WHITE, bullet_color=LIGHT_BLUE)
rect(s, Inches(4.43), Inches(6.3), Inches(3.3), Inches(0.5), BRIGHT_BLUE)
txb(s, "40 tokens → auto-completed",
    Inches(4.48), Inches(6.36), Inches(3.2), Inches(0.38),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(8.08), Inches(1.55), Inches(4.77), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ Screenshot ]\nAgreement Generator\n(RCPS / SAFE)",
    Inches(8.08), Inches(3.7), Inches(4.77), Inches(1.4),
    size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FEATURE 5: INVESTMENT EXECUTION TRACKER
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Feature 5  —  Investment Execution Tracker", "투자 집행 Tab")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), NAVY)
txb(s, "Investment Execution Tab",
    Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.55),
    size=18, bold=True, color=WHITE)
txb(s, "Track required steps after investment execution",
    Inches(0.7), Inches(2.28), Inches(5.6), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "Bank payment instruction (운용지시) tracking",
    "Post-payment document collection management",
    "Custodian bank 30-day deadline — hard cutoff, highlighted prominently",
    "Fund type:  Government-backed (모태)  vs  Private fund",
    "Investment structure (method · amount · equity %) recorded",
    "Email draft auto-generation (AI-assisted, for outgoing use)",
    "Auto-save — data preserved when navigating away",
], Inches(0.7), Inches(2.75), Inches(5.6), Inches(3.7),
   size=14, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(0.7), Inches(6.55), Inches(5.4), Inches(0.32), RED)
txb(s, "  Bank 30-day deadline — auto-surfaced at top of Action Center",
    Inches(0.82), Inches(6.58), Inches(5.1), Inches(0.26),
    size=12, bold=True, color=WHITE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.1), Inches(5.5), LIGHT_BLUE,
     line_color=MID_BLUE, line_width=1.5)
txb(s, "[ Screenshot ]\nInvestment Execution Tab",
    Inches(6.75), Inches(3.95), Inches(6.1), Inches(1.1),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FEATURE 6: SAFE CONVERSION + AI CHATBOT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Feature 6  —  SAFE Conversion Tracker  +  AI Process Assistant")

rect(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.55), NAVY)
txb(s, "SAFE Conversion Tracker Tab",
    Inches(0.7), Inches(1.7), Inches(5.7), Inches(0.55),
    size=17, bold=True, color=WHITE)
txb(s, "Step-by-step tracking of SAFE → equity conversion",
    Inches(0.7), Inches(2.28), Inches(5.7), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "3 stages:  Pre-conversion  ·  During conversion  ·  Post-conversion",
    "Deadline auto-calculation and D-day countdown display",
    "AI-assisted conversion share count calculation (Gemini)\n  — based on investment principal · Valuation Cap · discount rate",
    "Post-conversion registry change tracking",
    "Step-by-step checklist with per-item completion display",
], Inches(0.7), Inches(2.75), Inches(5.7), Inches(3.65),
   size=13.5, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(6.85), Inches(1.55), Inches(6.0), Inches(5.55), MID_BLUE)
txb(s, "AI Process Assistant",
    Inches(7.05), Inches(1.7), Inches(5.6), Inches(0.55),
    size=17, bold=True, color=WHITE)
txb(s, "Floating chatbot — bottom-right of every page",
    Inches(7.05), Inches(2.28), Inches(5.6), Inches(0.38),
    size=13, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "Grounded in SparkLabs internal process documents",
    "Context-aware — knows which deal you're currently viewing",
    "Answers in both Korean and English",
    "AI field suggestions for agreements (Gemini)",
    "Auto-extracts execution-related numbers (Gemini)",
    "Instant process guidance — no need to interrupt the mentor",
], Inches(7.05), Inches(2.75), Inches(5.6), Inches(3.65),
   size=13.5, color=WHITE, bullet_color=LIGHT_BLUE)

rect(s, Inches(0.5), Inches(7.13), Inches(12.35), Inches(0.35), DARK)
txb(s, "Full KO/EN bilingual toggle — entire platform switches between Korean and English instantly",
    Inches(0.65), Inches(7.17), Inches(12.0), Inches(0.28),
    size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AI TOOLS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "AI Tools Used", "Two Different Roles — In the App vs. Building the App")

rect(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.55), NAVY)
txb(s, "AI inside the app",
    Inches(0.7), Inches(1.72), Inches(5.7), Inches(0.38),
    size=12, bold=True, color=BRIGHT_BLUE)
txb(s, "Google Gemini API",
    Inches(0.7), Inches(2.13), Inches(5.7), Inches(0.6),
    size=22, bold=True, color=WHITE)
txb(s, "gemini-flash-lite-latest",
    Inches(0.7), Inches(2.75), Inches(5.7), Inches(0.38),
    size=12, color=LIGHT_BLUE, italic=True)
bullet_frame(s, [
    "Filename auto-classification  (Doc Collection tab)",
    "Diligence AI analysis — per-item review point suggestions",
    "Agreement field AI suggestions  (Agreement tab)",
    "Execution number auto-extraction  (Execution tab)",
    "SAFE conversion share count calculation  (Conversion tab)",
    "In-app chatbot assistant — process Q&A",
], Inches(0.7), Inches(3.2), Inches(5.6), Inches(3.55),
   size=13.5, color=WHITE, bullet_color=BRIGHT_BLUE)

rect(s, Inches(6.85), Inches(1.55), Inches(6.0), Inches(5.55), MID_BLUE)
txb(s, "AI used to build the app",
    Inches(7.05), Inches(1.72), Inches(5.6), Inches(0.38),
    size=12, bold=True, color=LIGHT_BLUE)

rect(s, Inches(7.05), Inches(2.15), Inches(5.6), Inches(1.85), NAVY)
txb(s, "Claude  (Anthropic)",
    Inches(7.25), Inches(2.27), Inches(5.2), Inches(0.5),
    size=17, bold=True, color=WHITE)
txb(s, "Code writing, debugging, feature design\nComplex DOCX XML manipulation",
    Inches(7.25), Inches(2.77), Inches(5.2), Inches(1.0),
    size=13, color=LIGHT_BLUE)

rect(s, Inches(7.05), Inches(4.13), Inches(5.6), Inches(1.85), NAVY)
txb(s, "Claude Code  (CLI tool)",
    Inches(7.25), Inches(4.25), Inches(5.2), Inches(0.5),
    size=17, bold=True, color=WHITE)
txb(s, "AI coding assistant in the terminal\nMade development possible with no prior coding experience",
    Inches(7.25), Inches(4.75), Inches(5.2), Inches(1.0),
    size=13, color=LIGHT_BLUE)

rect(s, 0, Inches(7.1), W, Inches(0.4), LIGHT_BLUE)
txb(s, "Full stack:   Next.js 16  ·  React 19  ·  TypeScript  ·  Tailwind CSS v4  ·  Vercel  ·  Vercel Blob  ·  Google Drive API  ·  Google Gemini API  ·  NextAuth",
    Inches(0.7), Inches(7.13), Inches(12.5), Inches(0.35),
    size=12, color=NAVY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — GEMINI API COST ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Google Gemini API — Cost Analysis", "How much does the AI actually cost?")

rect(s, Inches(0.5), Inches(1.55), Inches(3.7), Inches(5.55), NAVY)
txb(s, "Model used",
    Inches(0.7), Inches(1.72), Inches(3.3), Inches(0.38),
    size=12, bold=True, color=BRIGHT_BLUE)
txb(s, "gemini-flash-\nlite-latest",
    Inches(0.7), Inches(2.12), Inches(3.3), Inches(1.1),
    size=20, bold=True, color=WHITE)
txb(s, "Google's cheapest\ncurrent AI model",
    Inches(0.7), Inches(3.22), Inches(3.3), Inches(0.65),
    size=12, color=LIGHT_BLUE)
rect(s, Inches(0.7), Inches(4.0), Inches(3.1), Inches(0.06), BRIGHT_BLUE)
txb(s, "Cost per call",
    Inches(0.7), Inches(4.15), Inches(3.3), Inches(0.38),
    size=12, bold=True, color=BRIGHT_BLUE)
txb(s, "~$0.0002",
    Inches(0.7), Inches(4.55), Inches(3.3), Inches(0.75),
    size=26, bold=True, color=WHITE)
txb(s, "per API call",
    Inches(0.7), Inches(5.3), Inches(3.3), Inches(0.38),
    size=12, color=LIGHT_BLUE, italic=True)
txb(s, "13x cheaper than Claude Haiku",
    Inches(0.7), Inches(5.78), Inches(3.3), Inches(0.38),
    size=11, color=GRAY)

rect(s, Inches(4.45), Inches(1.55), Inches(5.3), Inches(5.55), WHITE,
     line_color=LIGHT_BLUE, line_width=1)
txb(s, "Estimated annual usage",
    Inches(4.65), Inches(1.72), Inches(4.9), Inches(0.42),
    size=14, bold=True, color=NAVY)

rows = [
    ("Filename classification",  "~150 calls", "~$0.03"),
    ("Diligence AI analysis",    "~50 calls",  "~$0.01"),
    ("Chatbot assistant",        "~250 calls", "~$0.05"),
]
row_cols = [LIGHT_BLUE, OFF_WHITE, LIGHT_BLUE]
for i, (label, calls, cost) in enumerate(rows):
    ty = Inches(2.28) + i * Inches(0.77)
    rect(s, Inches(4.55), ty, Inches(5.1), Inches(0.65), row_cols[i])
    txb(s, label, Inches(4.72), ty + Inches(0.12), Inches(2.4), Inches(0.4),
        size=13, color=DARK)
    txb(s, calls, Inches(7.2), ty + Inches(0.12), Inches(1.0), Inches(0.4),
        size=13, color=GRAY, align=PP_ALIGN.CENTER)
    txb(s, cost, Inches(8.25), ty + Inches(0.12), Inches(1.2), Inches(0.4),
        size=13, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

rect(s, Inches(4.55), Inches(4.59), Inches(5.1), Inches(0.82), NAVY)
txb(s, "Annual total", Inches(4.72), Inches(4.72), Inches(2.4), Inches(0.5),
    size=14, bold=True, color=WHITE)
txb(s, "~450 calls", Inches(7.2), Inches(4.72), Inches(1.0), Inches(0.5),
    size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
txb(s, "~$0.09", Inches(8.05), Inches(4.69), Inches(1.4), Inches(0.55),
    size=20, bold=True, color=BRIGHT_BLUE, align=PP_ALIGN.RIGHT)

txb(s, "* At 10x usage: ~$0.90/year   *  At 100x usage: ~$9.00/year",
    Inches(4.65), Inches(5.55), Inches(5.0), Inches(0.38),
    size=11, color=GRAY, italic=True)
txb(s, "Free tier: up to 60 requests/minute at no charge\n(current usage is comfortably within the free tier)",
    Inches(4.65), Inches(6.0), Inches(5.0), Inches(0.65),
    size=12, color=DARK)

rect(s, Inches(10.0), Inches(1.55), Inches(2.85), Inches(5.55), GREEN)
txb(s, "Annual AI\ncost",
    Inches(10.1), Inches(1.72), Inches(2.6), Inches(0.85),
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "$0.09",
    Inches(10.1), Inches(2.55), Inches(2.6), Inches(1.5),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "per year",
    Inches(10.1), Inches(4.05), Inches(2.6), Inches(0.45),
    size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
rect(s, Inches(10.2), Inches(4.6), Inches(2.4), Inches(0.05), WHITE)
txb(s, "Starbucks\nAmericano",
    Inches(10.1), Inches(4.75), Inches(2.6), Inches(0.75),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "~$5.00",
    Inches(10.1), Inches(5.5), Inches(2.6), Inches(0.65),
    size=20, bold=True, color=rgb(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)
txb(s, "1 coffee > 1 year of AI",
    Inches(10.1), Inches(6.2), Inches(2.6), Inches(0.45),
    size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — RESULTS & IMPROVEMENTS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, NAVY)
header_bar(s, "Results & Improvements")
rect(s, 0, Inches(1.3), W, Inches(0.06), BRIGHT_BLUE)

rect(s, Inches(0.5), Inches(1.6), Inches(5.5), Inches(3.75), RED)
txb(s, "Before",
    Inches(0.7), Inches(1.75), Inches(5.1), Inches(0.52),
    size=17, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "2~3\nhours",
    Inches(0.7), Inches(2.28), Inches(5.1), Inches(1.5),
    size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "per agreement\n(manual input + review)",
    Inches(0.7), Inches(3.82), Inches(5.1), Inches(0.7),
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "↓  Every blank filled by hand",
    Inches(0.7), Inches(4.65), Inches(5.1), Inches(0.48),
    size=12, color=rgb(255, 200, 200), align=PP_ALIGN.CENTER)

txb(s, "→", Inches(5.5), Inches(2.9), Inches(0.9), Inches(1.5),
    size=50, bold=True, color=BRIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(6.3), Inches(1.6), Inches(6.55), Inches(3.75), GREEN)
txb(s, "After",
    Inches(6.5), Inches(1.75), Inches(6.2), Inches(0.52),
    size=17, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "~3\nmin",
    Inches(6.5), Inches(2.28), Inches(6.2), Inches(1.5),
    size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "per agreement\n(fill form → download)",
    Inches(6.5), Inches(3.82), Inches(6.2), Inches(0.7),
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "↓  69 fields auto-completed",
    Inches(6.5), Inches(4.65), Inches(6.2), Inches(0.48),
    size=12, color=rgb(200, 255, 200), align=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(5.52), Inches(12.35), Inches(0.48), MID_BLUE)
txb(s, "More:  auto bank deadline alerts  ·  centralized doc tracking  ·  structured diligence  ·  instant AI answers  ·  bilingual support",
    Inches(0.65), Inches(5.6), Inches(12.0), Inches(0.35),
    size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(6.12), Inches(12.35), Inches(1.0), BRIGHT_BLUE)
txb(s, "Over 97% reduction in time per agreement  (2-3 hours → ~3 minutes)",
    Inches(0.7), Inches(6.18), Inches(12.0), Inches(0.82),
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — LIMITATIONS & ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, OFF_WHITE)
header_bar(s, "Limitations & Future Roadmap")

rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5), NAVY)
txb(s, "Current Limitations",
    Inches(0.7), Inches(1.72), Inches(5.5), Inches(0.58),
    size=20, bold=True, color=ORANGE)
bullet_frame(s, [
    "Agreement template updates are manual\n  (XML must be edited directly when legal clauses change)",
    "Multiple interested parties not yet supported (currently 1 per deal)",
    "Not fully optimized for mobile",
    "No e-signature — wet signatures still required",
    "No user permission levels — all team members see everything",
    "Diligence checklist item count is hard to adjust",
], Inches(0.7), Inches(2.45), Inches(5.5), Inches(4.4),
   size=14, color=WHITE, bullet_color=ORANGE)

rect(s, Inches(6.8), Inches(1.55), Inches(6.05), Inches(5.5), LIGHT_BLUE)
txb(s, "Planned Features",
    Inches(7.0), Inches(1.72), Inches(5.6), Inches(0.58),
    size=20, bold=True, color=NAVY)
bullet_frame(s, [
    "Email notifications (contract sent / signed)",
    "E-signature integration  (e.g. DocuSign)",
    "AI-powered contract clause review",
    "Fund-level investment statistics dashboard",
    "Multiple interested party support",
    "Contract version control (track revision history)",
    "Admin UI for template editing",
], Inches(7.0), Inches(2.45), Inches(5.6), Inches(4.4),
   size=14, color=DARK, bullet_color=MID_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, NAVY)

rect(s, 0, H - Inches(0.55), W, Inches(0.55), BRIGHT_BLUE)
rect(s, 0, Inches(0),        W, Inches(0.55), BRIGHT_BLUE)

txb(s, "Live Demo",
    Inches(1.5), Inches(1.8), Inches(10), Inches(2.0),
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Live Demonstration of the Website",
    Inches(1.5), Inches(3.82), Inches(10), Inches(0.75),
    size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(4.78), Inches(6.35), Inches(0.08), BRIGHT_BLUE)

txb(s, "Demonstrating the actual site live",
    Inches(1.5), Inches(5.02), Inches(10), Inches(0.6),
    size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

demo_items = ["Dashboard view", "Document Collection tab", "Agreement drafting & download", "AI chatbot demo"]
txb(s, "     ·     ".join(demo_items),
    Inches(0.5), Inches(5.88), Inches(12.35), Inches(0.55),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "sparklabs-presentation-english-v2.pptx")
prs.save(out)
print(f"Saved: {out}")
